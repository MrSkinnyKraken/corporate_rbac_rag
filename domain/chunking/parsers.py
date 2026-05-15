"""Multi-format document parser dispatcher.

Extracts plain text from a file path based on its extension, returning a
single string ready for the :class:`~domain.chunking.core_chunker.CustomRBACChunker`.
The parser is purely deterministic — no PII detection, no chunking, no LLM —
and is entirely independent of the infrastructure adapters.

Supported formats
-----------------

================  ===============================  =====================================
Extension(s)      Library                          Notes
================  ===============================  =====================================
``.txt``,         stdlib                           Read as UTF-8 plain text.
``.md``
``.csv``          stdlib (``csv``)                 Rendered as pipe-separated rows so the
                                                   chunker treats each row as one line —
                                                   the interval-merge fix from Phase 1.
``.pdf``          ``pypdf``                        Page text concatenated; encrypted PDFs
                                                   raise :class:`DocumentProcessingError`.
``.docx``         ``python-docx``                  Paragraphs only; tables and footers
                                                   are NOT extracted in this baseline.
``.xlsx``         ``openpyxl``                     Every sheet is rendered as
                                                   ``"Sheet: NAME"`` + pipe-separated rows.
                                                   The ``Sheet:`` prefix matches the
                                                   chunker's structural separator.
``.pptx``         ``python-pptx``                  Slide titles + text frames; embedded
                                                   charts and SmartArt are skipped.
``.html``,        ``beautifulsoup4``               ``<script>``, ``<style>``, ``<noscript>``
``.htm``                                           tags are stripped before extraction.
================  ===============================  =====================================

Anything else raises :class:`~core.exceptions.DocumentProcessingError` with a
clear diagnostic message and the offending source path. Unexpected exceptions
inside a registered handler are likewise wrapped, so callers only have to
catch the project's :class:`~core.exceptions.PipelineError` hierarchy.

Future formats (``.xls``, ``.rtf``, ``.odt``, ``.epub``, ``.json``, ``.xml``)
can be added by writing one private ``_parse_<ext>`` static method and a
single dispatch entry; the rest of the chunking pipeline does not change.
"""

from __future__ import annotations

import csv
import io
from pathlib import Path
from typing import Callable

from bs4 import BeautifulSoup
from docx import Document as _DocxDocument
from openpyxl import load_workbook
from pptx import Presentation as _PptxPresentation
from pypdf import PdfReader

from core.exceptions import DocumentProcessingError


# Type alias: a parser handler receives an absolute path and returns text.
_Handler = Callable[[Path], str]


class DocumentParser:
    """Extension-based dispatcher that returns plain text for a given file.

    The dispatch table is built at construction time so handlers are bound
    static methods of this instance. Subclasses can override individual
    ``_parse_*`` methods (e.g., to swap ``pypdf`` for ``pdfplumber``) without
    touching the public :meth:`parse` API.
    """

    def __init__(self) -> None:
        self._dispatch: dict[str, _Handler] = {
            # Plain-text family
            ".txt":  self._parse_text,
            ".md":   self._parse_text,
            # Tabular
            ".csv":  self._parse_csv,
            ".xlsx": self._parse_xlsx,
            # Office documents
            ".pdf":  self._parse_pdf,
            ".docx": self._parse_docx,
            ".pptx": self._parse_pptx,
            # Web
            ".html": self._parse_html,
            ".htm":  self._parse_html,
        }

    # ───────────────────────────── Public API ────────────────────────────────
    def parse(self, path: Path) -> str:
        """Extract plain text from ``path``.

        Args:
            path: Absolute or relative path to the document. Must exist.

        Returns:
            The document's textual content as a single ``str`` (with native
            line breaks preserved).

        Raises:
            DocumentProcessingError: if the file does not exist, the
                extension is not registered, or the underlying parser
                library raises any error during extraction.
        """
        if not path.exists():
            raise DocumentProcessingError(
                f"File does not exist: {path}",
                source=str(path),
            )
        if not path.is_file():
            raise DocumentProcessingError(
                f"Path is not a regular file: {path}",
                source=str(path),
            )

        ext: str = path.suffix.lower()
        handler: _Handler | None = self._dispatch.get(ext)
        if handler is None:
            raise DocumentProcessingError(
                f"Unsupported extension {ext!r} for {path.name}; "
                f"supported extensions: {sorted(self._dispatch.keys())}.",
                source=str(path),
            )

        try:
            return handler(path)
        except DocumentProcessingError:
            # Already a domain error — propagate without re-wrapping.
            raise
        except Exception as exc:  # noqa: BLE001  (intentional translation boundary)
            raise DocumentProcessingError(
                f"Failed to parse {path.name} with handler {handler.__name__!r}: "
                f"{type(exc).__name__}: {exc}",
                source=str(path),
                details={"handler": handler.__name__, "extension": ext},
            ) from exc

    @property
    def supported_extensions(self) -> list[str]:
        """Return the list of registered file extensions, sorted."""
        return sorted(self._dispatch.keys())

    # ─────────────────────── Plain-text format handlers ──────────────────────
    @staticmethod
    def _parse_text(path: Path) -> str:
        """Read ``.txt`` or ``.md`` as UTF-8 with replacement on bad bytes."""
        return path.read_text(encoding="utf-8", errors="replace")

    # ──────────────────────────── Tabular formats ────────────────────────────
    @staticmethod
    def _parse_csv(path: Path) -> str:
        """Render CSV as pipe-separated rows.

        Each row becomes one line of the form ``"col1 | col2 | col3"``.
        Keeping every row on a single line is what enables the chunker's
        line-extension + interval-merge logic to treat a tabular row with
        multiple PII matches as a *single* sensitive fragment, rather than
        the over-fragmented mess Phase 1 produced on the
        ``clients-and-billings.xlsx`` corpus.
        """
        out = io.StringIO()
        with path.open("r", encoding="utf-8", newline="") as fp:
            reader = csv.reader(fp)
            for row in reader:
                # Strip cells but keep empty placeholders to preserve column
                # alignment (downstream regex is column-agnostic anyway).
                out.write(" | ".join(cell.strip() for cell in row) + "\n")
        return out.getvalue()

    @staticmethod
    def _parse_xlsx(path: Path) -> str:
        """Render every sheet of an ``.xlsx`` workbook as pipe-separated rows.

        Each sheet is preceded by ``"Sheet: NAME"`` — the prefix that the
        chunker's :data:`STRUCTURAL_SEPARATOR_PATTERNS` already recognises,
        so a multi-sheet workbook splits cleanly into one chunk per sheet
        before per-row PII detection runs.

        The workbook is opened with ``read_only=True`` (faster, lower memory
        on large files) and ``data_only=True`` (returns calculated values
        instead of formula text — matches what a human would see).
        """
        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            sections: list[str] = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sections.append(f"Sheet: {sheet_name}")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cell.strip() for cell in cells):
                        sections.append(" | ".join(cells))
                sections.append("")  # blank line between sheets
            return "\n".join(sections)
        finally:
            workbook.close()

    # ──────────────────────────── Office documents ───────────────────────────
    @staticmethod
    def _parse_pdf(path: Path) -> str:
        """Concatenate the text of every PDF page.

        Pages are separated by a blank line so the structural splitter in
        the chunker can use ``\\n\\n`` as a paragraph boundary on multi-page
        documents.
        """
        reader = PdfReader(str(path))
        if reader.is_encrypted:
            raise DocumentProcessingError(
                f"PDF is encrypted and cannot be read: {path.name}",
                source=str(path),
                details={"pages": len(reader.pages)},
            )
        pages: list[str] = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            pages.append(page_text)
        return "\n\n".join(pages)

    @staticmethod
    def _parse_docx(path: Path) -> str:
        """Concatenate every paragraph of a Word document.

        Tables and headers/footers are NOT extracted in this baseline
        implementation. If a future ingestion needs them, override this
        method or replace it with a richer adapter (e.g., ``unstructured``).
        """
        document = _DocxDocument(str(path))
        return "\n".join(para.text for para in document.paragraphs)

    @staticmethod
    def _parse_pptx(path: Path) -> str:
        """Concatenate the text frames of every slide.

        Each slide is preceded by a ``"Slide N:"`` marker so the chunker's
        structural splitter can isolate slide content. Only paragraph text
        inside ``text_frame`` shapes is extracted; embedded charts, SmartArt
        diagrams, and image OCR are skipped — those would need an OCR
        pipeline (Tesseract, Vision LLM, …) outside this module's scope.
        """
        presentation = _PptxPresentation(str(path))
        sections: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            sections.append(f"Slide {index}:")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        sections.append(text)
            sections.append("")  # blank line between slides
        return "\n".join(sections)

    # ────────────────────────────── Web formats ──────────────────────────────
    @staticmethod
    def _parse_html(path: Path) -> str:
        """Strip HTML tags and return the visible text.

        Removes ``<script>``, ``<style>``, and ``<noscript>`` blocks first
        to avoid polluting the chunker with code or CSS. Whitespace is
        normalised line-by-line; consecutive empty lines collapse, but
        paragraph structure is otherwise preserved so the structural
        splitter can find paragraph breaks.

        For HTML pages with heavy JavaScript-rendered content (SPAs),
        :mod:`beautifulsoup4` only sees the static markup. Such pages
        should be pre-rendered (Playwright, Selenium) before ingestion;
        that work belongs to the application layer, not this parser.
        """
        with path.open("r", encoding="utf-8", errors="replace") as fp:
            soup = BeautifulSoup(fp, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        raw_text: str = soup.get_text(separator="\n")
        # Collapse whitespace per line; drop fully-blank lines.
        lines: list[str] = [line.strip() for line in raw_text.splitlines()]
        return "\n".join(line for line in lines if line)
